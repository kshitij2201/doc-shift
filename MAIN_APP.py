from flask import Flask, request, send_file, jsonify, render_template, session, redirect, url_for
from functools import wraps
from PIL import Image
import os
import io
import uuid
import sqlite3
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from io import BytesIO
from pdf2image import convert_from_bytes
from docx import Document
from openpyxl import load_workbook
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from rembg import remove
import logging
import subprocess
import fitz  # PyMuPDF for AI PDF Editor
from fpdf import FPDF  # For AI PDF Editor
import glob
from flask_cors import CORS
import requests
import json
from collections import deque
from gtts import gTTS
import speech_recognition as sr
import tempfile
from pydub import AudioSegment
import time
from bs4 import BeautifulSoup

app = Flask(__name__)
CORS(app)  # Enable CORS for API requests
app.secret_key = os.urandom(24)

# Configure logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Folder configurations
UPLOAD_FOLDER = 'Uploads'
CONVERTED_FOLDER = 'converted'
OUTPUT_FOLDER = 'converted'
STATIC_FOLDER = 'static'
AUDIO_FOLDER = 'tts_audio'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(CONVERTED_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
os.makedirs(STATIC_FOLDER, exist_ok=True)
os.makedirs(AUDIO_FOLDER, exist_ok=True)

# OpenRouter API configuration
OPENROUTER_API_KEY = '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'
OPENROUTER_API_URL = 'https://openrouter.ai/api/v1/chat/completions'
OPENROUTER_MODEL = 'gpt-4o-mini'

# Initialize speech recognizer for Speech to Text
recognizer = sr.Recognizer()

# Document Screener global variables
current_document_text = ''
conversation_history = deque(maxlen=10)  # Limit to last 10 messages for context

# AI PDF Editor global variable
latest_text = ""

# Clean up old PDFs in UPLOAD_FOLDER
for f in glob.glob(os.path.join(UPLOAD_FOLDER, "*.pdf")):
    os.remove(f)

def init_db():
    conn = sqlite3.connect('file_conversion.db')
    c = conn.cursor()
    c.execute('''
        CREATE TABLE IF NOT EXISTS conversions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            conversion_type TEXT NOT NULL,
            original_filename TEXT NOT NULL,
            converted_filename TEXT NOT NULL,
            timestamp DATETIME DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    try:
        c.execute('ALTER TABLE conversions ADD COLUMN file_path TEXT')
    except sqlite3.OperationalError:
        pass
    c.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT NOT NULL UNIQUE,
            password_hash TEXT NOT NULL
        )
    ''')
    try:
        default_username = 'admin'
        default_password = 'password'
        password_hash = generate_password_hash(default_password)
        c.execute('INSERT INTO users (username, password_hash) VALUES (?, ?)', (default_username, password_hash))
        conn.commit()
    except sqlite3.IntegrityError:
        pass
    conn.close()

init_db()

def log_conversion(conversion_type, original_filename, converted_filename, file_path=None):
    conn = sqlite3.connect('file_conversion.db')
    c = conn.cursor()
    c.execute('''
        INSERT INTO conversions (conversion_type, original_filename, converted_filename, file_path)
        VALUES (?, ?, ?, ?)
    ''', (conversion_type, original_filename, converted_filename, file_path))
    conn.commit()
    conn.close()

def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'logged_in' not in session:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

def is_ghostscript_installed():
    try:
        subprocess.run(['gs', '--version'], stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
        return True
    except (subprocess.CalledProcessError, FileNotFoundError):
        return False

# Document Screener Helper Functions
def extract_text_from_pdf(file_path):
    try:
        with open(file_path, 'rb') as file:
            reader = PdfReader(file)
            text = ''
            for page in reader.pages:
                text += page.extract_text() or ''
            return text
    except Exception as e:
        return f"Error extracting text from PDF: {str(e)}"

def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        text = '\n'.join([para.text for para in doc.paragraphs if para.text])
        return text
    except Exception as e:
        return f"Error extracting text from DOCX: {str(e)}"

def extract_text_from_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        return f"Error extracting text from TXT: {str(e)}"

def analyze_text_with_openrouter(text, format_type):
    prompt = (
        f"Analyze the following document text and provide a summary in {'a concise paragraph' if format_type == 'paragraph' else 'bullet points'}. "
        f"Focus on key themes, topics, or entities mentioned in the text.\n\nText:\n{text[:2000]}"
    )
    
    headers = {
        'Authorization': f'Bearer {OPENROUTER_API_KEY}',
        'Content-Type': 'application/json'
    }
    
    data = {
        'model': 'mistralai/mixtral-8x7b-instruct',
        'messages': [{'role': 'user', 'content': prompt}]
    }
    
    try:
        response = requests.post(OPENROUTER_API_URL, headers=headers, json=data)
        response.raise_for_status()
        result = response.json()
        return result['choices'][0]['message']['content']
    except Exception as e:
        logger.error(f"OpenRouter API error: {str(e)}")
        return f"Error analyzing text with OpenRouter: {str(e)}"

def chat_with_openrouter(message):
    global current_document_text
    context = f"Document text:\n\n{current_document_text[:2000]}\n\nConversation history:\n"
    for role, msg in conversation_history:
        context += f"{role}: {msg}\n"
    prompt = f"{context}\nUser: {message}\nAssistant: Answer based on the document and conversation history. If the question is about names or specific details, extract relevant information from the document. If no relevant information is found, say so clearly."

    headers = {
        'Authorization': f'Bearer {OPENROUTER_API_KEY}',
        'Content-Type': 'application/json'
    }
    
    data = {
        'model': 'mistralai/mixtral-8x7b-instruct',
        'messages': [{'role': 'user', 'content': prompt}]
    }
    
    try:
        response = requests.post(OPENROUTER_API_URL, headers=headers, json=data)
        response.raise_for_status()
        result = response.json()
        response_text = result.get('choices', [{}])[0].get('message', {}).get('content', '')
        conversation_history.append(('User', message))
        conversation_history.append(('Assistant', response_text))
        return response_text
    except Exception as e:
        logger.error(f"OpenRouter API error: {str(e)}")
        return f"Error chatting with OpenRouter: {str(e)}"

# Plagiarism Scanner Helper Functions
def fetch_web_snippets(query, max_results=5):
    """Scrape DuckDuckGo search results for snippets."""
    search_url = f"https://html.duckduckgo.com/html/?q={query}"
    headers = {
        "User-Agent": "Mozilla/5.0"
    }
    try:
        response = requests.get(search_url, headers=headers, timeout=5)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        results = soup.find_all('a', {'class': 'result__a'}, limit=max_results)
        snippets = []
        for r in results:
            text = r.get_text(strip=True)
            if text:
                snippets.append(text)
        return snippets
    except Exception as e:
        logger.error(f"Error fetching web snippets: {str(e)}")
        return []

def call_openrouter_similarity(text_a, text_b):
    """Use OpenRouter AI to compare similarity."""
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json",
        "User-Agent": "Mozilla/5.0"
    }
    messages = [
        {
            "role": "system",
            "content": "You are a plagiarism detection assistant. Given two texts, respond ONLY with a plagiarism similarity percentage (0 to 100) and a brief explanation, separated by a newline."
        },
        {
            "role": "user",
            "content": f"Text A:\n{text_a}\n\nText B:\n{text_b}"
        }
    ]
    data = {
        "model": OPENROUTER_MODEL,
        "messages": messages
    }
    try:
        response = requests.post(OPENROUTER_API_URL, headers=headers, json=data)
        response.raise_for_status()
        return response.json()['choices'][0]['message']['content'].strip()
    except Exception as e:
        logger.error(f"OpenRouter API error in similarity check: {str(e)}")
        return f"Error: {str(e)}"

# AI PDF Editor Helper Function
def extract_structured_text(path):
    try:
        doc = fitz.open(path)
        structured_data = []
        for page in doc:
            blocks = page.get_text("dict")['blocks']
            for block in blocks:
                if 'lines' in block:
                    for line in block['lines']:
                        line_text = " ".join([span['text'] for span in line['spans']])
                        structured_data.append(line_text)
        doc.close()
        return "\n".join(structured_data)
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        return f"Error extracting text from PDF: {str(e)}"

def retry_api_request(url, headers, data, max_retries=3, delay=2):
    """Retry API request with exponential backoff."""
    for attempt in range(max_retries):
        try:
            response = requests.post(url, headers=headers, json=data, timeout=30)
            logger.debug(f"API response status: {response.status_code}")
            if response.status_code == 429:
                logger.warning(f"Rate limit hit on attempt {attempt + 1}, retrying after {delay} seconds")
                time.sleep(delay)
                delay *= 2  # Exponential backoff
                continue
            response.raise_for_status()
            return response
        except requests.exceptions.RequestException as e:
            logger.error(f"API request failed on attempt {attempt + 1}: {str(e)}")
            if attempt == max_retries - 1:
                raise e
            time.sleep(delay)
            delay *= 2
    raise Exception("Max retries exceeded for API request")

@app.route('/login', methods=['GET', 'POST'])
def login():
    if 'logged_in' in session:
        return redirect(url_for('index'))
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        conn = sqlite3.connect('file_conversion.db')
        c = conn.cursor()
        c.execute('SELECT password_hash FROM users WHERE username = ?', (username,))
        user = c.fetchone()
        conn.close()
        if user and check_password_hash(user[0], password):
            session['logged_in'] = True
            return redirect(url_for('index'))
        else:
            return render_template('login.html', error='Invalid username or password')
    return render_template('login.html', error=None)

@app.route('/logout')
@login_required
def logout():
    session.pop('logged_in', None)
    return redirect(url_for('login'))

@app.route('/')
@login_required
def index():
    return render_template('index.html')

@app.route('/image-to-pdf')
@login_required
def image_to_pdf_page():
    return render_template('image_to_pdf.html')

@app.route('/pdf-to-image')
@login_required
def pdf_to_image_page():
    return render_template('pdf_to_image.html')

@app.route('/merge-pdfs')
@login_required
def merge_pdfs_page():
    return render_template('merge_pdfs.html')

@app.route('/word-to-pdf')
@login_required
def word_to_pdf_page():
    return render_template('word_to_pdf.html')

@app.route('/excel-to-pdf')
@login_required
def excel_to_pdf_page():
    return render_template('excel_to_pdf.html')

@app.route('/pdf-to-ppt')
@login_required
def pdf_to_ppt_page():
    return render_template('pdf_to_ppt.html')

@app.route('/bg-remover')
@login_required
def bg_remover_page():
    return render_template('bg_remover.html')

@app.route('/admin-logs')
@login_required
def logs_page():
    return render_template('logs.html')

@app.route('/compress-pdf')
@login_required
def compress_pdf_page():
    return render_template('compress_pdf.html')

@app.route('/split-pdf')
@login_required
def split_pdf_page():
    return render_template('split_pdf.html')

@app.route('/remove-pages-ui')
@login_required
def remove_pages_ui():
    return render_template('remove_page.html')

@app.route('/document-screener')
@login_required
def document_screener_page():
    global current_document_text, conversation_history
    current_document_text = ''
    conversation_history.clear()
    return render_template('document_screener.html')

@app.route('/plagiarism-scanner')
@login_required
def plagiarism_scanner_page():
    result = session.pop('plagiarism_result', None)
    input_text = session.pop('plagiarism_input_text', '')
    return render_template('plagiarism.html', result=result, input_text=input_text)

@app.route('/text-to-speech')
@login_required
def text_to_speech_page():
    return render_template('text_to_speech.html')

@app.route('/speech-to-text')
@login_required
def speech_to_text_page():
    return render_template('speech_to_text.html')

@app.route('/ai-pdf-editor')
@login_required
def ai_pdf_editor_page():
    return render_template('ai_pdf_editor.html')

@app.route('/text-summarizer')
@login_required
def text_summarizer_page():
    return render_template('text_summarizer.html')

@app.route('/convert/image-to-pdf', methods=['POST'])
@login_required
def convert_image_to_pdf():
    if 'images' not in request.files:
        return jsonify({'error': 'No images provided'}), 400

    files = request.files.getlist('images')
    image_list = []

    for file in files:
        try:
            image = Image.open(file.stream)
            if image.mode != 'RGB':
                image = image.convert('RGB')
            image_list.append(image)
        except Exception as e:
            return jsonify({'error': f'Failed to read image: {str(e)}'}), 500

    if not image_list:
        return jsonify({'error': 'No valid images found'}), 400

    output_filename = f"{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(CONVERTED_FOLDER, output_filename)

    try:
        image_list[0].save(output_path, save_all=True, append_images=image_list[1:], format='PDF')
        log_conversion('image-to-pdf', files[0].filename, output_filename, output_path)
    except Exception as e:
        return jsonify({'error': f'PDF conversion failed: {str(e)}'}), 500

    return send_file(output_path, as_attachment=True, download_name="converted.pdf", mimetype='application/pdf')

@app.route('/convert/pdf-to-image', methods=['POST'])
@login_required
def convert_pdf_to_image():
    pdf_file = request.files.get('pdf')
    if not pdf_file:
        return "No PDF uploaded", 400

    try:
        images = convert_from_bytes(pdf_file.read(), fmt='png', single_file=True)
        output_filename = f"{uuid.uuid4().hex}.png"
        output_path = os.path.join(CONVERTED_FOLDER, output_filename)
        images[0].save(output_path, format='PNG')

        log_conversion('pdf-to-image', pdf_file.filename, output_filename, output_path)

        return send_file(output_path, as_attachment=True, download_name="converted.png", mimetype='image/png')
    except Exception as e:
        return f"Error during conversion: {str(e)}", 500

@app.route('/merge/pdfs', methods=['POST'])
@login_required
def merge_pdfs():
    if 'pdfs' not in request.files:
        return jsonify({'error': 'No PDF files provided'}), 400

    files = request.files.getlist('pdfs')
    merger = PdfMerger()

    try:
        for file in files:
            merger.append(file)

        output_filename = f"merged_{uuid.uuid4().hex}.pdf"
        output_path = os.path.join(CONVERTED_FOLDER, output_filename)
        merger.write(output_path)
        merger.close()

        log_conversion('merge-pdfs', files[0].filename, output_filename, output_path)
        return send_file(output_path, as_attachment=True, download_name='merged.pdf', mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Merge failed: {str(e)}'}), 500

@app.route('/admin/logs')
@login_required
def view_logs():
    conn = sqlite3.connect('file_conversion.db')
    c = conn.cursor()
    c.execute('SELECT * FROM conversions ORDER BY timestamp DESC')
    rows = c.fetchall()
    conn.close()
    return jsonify(rows)

@app.route('/download')
@login_required
def download():
    file_path = request.args.get('file_path')
    file_name = request.args.get('file_name')
    mime_type = request.args.get('mime_type', 'application/octet-stream')
    if not os.path.exists(file_path):
        return jsonify({'error': 'File not found on server'}), 404
    return send_file(file_path, as_attachment=True, download_name=file_name, mimetype=mime_type)

@app.route('/convert_word_to_pdf', methods=['POST'])
@login_required
def convert_word_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded!'}), 400

    file = request.files['file']
    if not file.filename.endswith('.docx'):
        return jsonify({'error': 'Only .docx files are supported'}), 400

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_filename = f"{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(OUTPUT_FOLDER, output_filename)

    try:
        file.save(input_path)
        doc = Document(input_path)
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        pdf_io = BytesIO()
        c = canvas.Canvas(pdf_io, pagesize=letter)
        y = 750
        for line in text:
            if y < 50:
                c.showPage()
                y = 750
            c.drawString(50, y, line[:100])
            y -= 15
        c.save()
        pdf_io.seek(0)
        log_conversion('word-to-pdf', file.filename, output_filename, output_path)
        with open(output_path, 'wb') as f:
            f.write(pdf_io.getvalue())
        return send_file(output_path, as_attachment=True, download_name=file.filename.replace('.docx', '.pdf'), mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Word to PDF conversion failed: {str(e)}'}), 500
    finally:
        if os.path.exists(input_path):
            os.remove(input_path)

@app.route('/convert_excel_to_pdf', methods=['POST'])
@login_required
def convert_excel_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded!'}), 400

    file = request.files['file']
    if not file.filename.endswith('.xlsx'):
        return jsonify({'error': 'Only .xlsx files are supported'}), 400

    input_path = os.path.join(UPLOAD_FOLDER, secure_filename(file.filename))
    output_filename = f"{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(OUTPUT_FOLDER, output_filename)

    try:
        file.save(input_path)
        wb = load_workbook(input_path)
        sheet = wb.active
        text = []
        for row in sheet.rows:
            row_text = [str(cell.value) if cell.value is not None else '' for cell in row]
            text.append(' | '.join(row_text))
        pdf_io = BytesIO()
        c = canvas.Canvas(pdf_io, pagesize=letter)
        y = 750
        for line in text:
            if y < 50:
                c.showPage()
                y = 750
            c.drawString(50, y, line[:100])
            y -= 15
        c.save()
        pdf_io.seek(0)
        log_conversion('excel-to-pdf', file.filename, output_filename, output_path)
        with open(output_path, 'wb') as f:
            f.write(pdf_io.getvalue())
        return send_file(output_path, as_attachment=True, download_name=file.filename.replace('.xlsx', '.pdf'), mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Excel to PDF conversion failed: {str(e)}'}), 500
    finally:
        if os.path.exists(input_path):
            os.remove(input_path)

@app.route('/convert_pdf_to_ppt', methods=['POST'])
@login_required
def convert_pdf_to_ppt():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded!'}), 400

    file = request.files['file']
    if not file.filename.endswith('.pdf'):
        return jsonify({'error': 'Only .pdf files are supported'}), 400

    output_filename = f"{uuid.uuid4().hex}.pptx"
    output_path = os.path.join(OUTPUT_FOLDER, output_filename)

    try:
        pdf_reader = PdfReader(file.stream)
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]

        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text = page.extract_text() or f"Page {page_num + 1}"
            lines = text.split('\n')
            filtered_lines = [line.strip() for line in lines if line.strip()]
            slide_text = '\n'.join(filtered_lines[:10])

            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = f"Slide {page_num + 1}"
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.text = slide_text
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        ppt_io = BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)

        log_conversion('pdf-to-ppt', file.filename, output_filename, output_path)
        with open(output_path, 'wb') as f:
            f.write(ppt_io.getvalue())
        return send_file(output_path, as_attachment=True, download_name=file.filename.replace('.pdf', '.pptx'), mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    except Exception as e:
        return jsonify({'error': f'PDF to PPT conversion failed: {str(e)}'}), 500

@app.route('/remove_background', methods=['POST'])
@login_required
def remove_background():
    logger.debug("Received request at /remove_background")
    
    if 'image' not in request.files:
        logger.error("No image uploaded")
        return jsonify({'error': 'No image uploaded!'}), 400

    file = request.files['image']
    valid_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp')
    if not file.filename.lower().endswith(valid_extensions):
        logger.error(f"Unsupported file extension: {file.filename}")
        return jsonify({'error': 'Only JPG, PNG, GIF, BMP, TIFF, and WEBP files are supported'}), 400

    if file.content_length > 10 * 1024 * 1024:
        logger.error("File size exceeds 10MB limit")
        return jsonify({'error': 'File size exceeds 10MB limit'}), 400

    try:
        logger.debug("Opening image with PIL")
        image = Image.open(file.stream)
        logger.debug(f"Image size: {image.size}")
        
        if image.size[0] > 4096 or image.size[1] > 4096:
            logger.error("Image resolution exceeds 4096x4096 limit")
            return jsonify({'error': 'Image resolution exceeds 4096x4096 limit'}), 400

        if file.filename.lower().endswith('.gif'):
            logger.debug("Processing GIF, using first frame")
            image.seek(0)

        if image.mode != 'RGBA':
            logger.debug("Converting image to RGBA mode")
            image = image.convert('RGBA')

        logger.debug("Starting background removal with rembg")
        output_image = remove(image)
        logger.debug("Background removal completed")

        output_filename = f"bg_removed_{uuid.uuid4().hex}.png"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        logger.debug(f"Saving output to disk: {output_path}")
        output_image.save(output_path, format='PNG')

        log_conversion('background-remover', file.filename, output_filename, output_path)

        logger.debug("Sending processed image back to client")
        return send_file(output_path, as_attachment=True, download_name='background_removed.png', mimetype='image/png')
    except Exception as e:
        logger.error(f"Background removal failed: {str(e)}", exc_info=True)
        return jsonify({'error': f'Background removal failed: {str(e)}'}), 500

@app.route('/get-page-count', methods=['POST'])
@login_required
def get_page_count():
    pdf = request.files.get('pdf')
    if not pdf:
        return jsonify({'error': 'No PDF uploaded'}), 400

    filename = secure_filename(pdf.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    pdf.save(filepath)

    try:
        reader = PdfReader(filepath)
        count = len(reader.pages)
        os.remove(filepath)
        return jsonify({'page_count': count})
    except Exception as e:
        os.remove(filepath)
        return jsonify({'error': str(e)}), 500

@app.route('/remove-pages', methods=['POST'])
@login_required
def remove_pages():
    pdf = request.files.get('pdf')
    removed_pages = request.form.get('removed_pages', '')
    total_pages = int(request.form.get('page_count', 0))

    if not pdf:
        return "No PDF uploaded", 400

    try:
        remove_indices = [int(i) for i in removed_pages.split(',') if i.strip().isdigit()]
    except ValueError:
        return "Invalid page indices", 400

    filename = secure_filename(pdf.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    pdf.save(filepath)

    reader = PdfReader(filepath)
    writer = PdfWriter()

    for i in range(len(reader.pages)):
        if i not in remove_indices:
            writer.add_page(reader.pages[i])

    output_filename = f"removed_{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(STATIC_FOLDER, output_filename)
    with open(output_path, 'wb') as f:
        writer.write(f)

    os.remove(filepath)
    log_conversion('remove-pages', filename, output_filename, output_path)
    return send_file(output_path, as_attachment=True, download_name='removed_pages.pdf')

@app.route('/compress', methods=['POST'])
@login_required
def compress_pdf():
    pdf_file = request.files.get('pdf')
    compression_level = request.form.get('compression_level')

    if not pdf_file:
        return "No PDF uploaded", 400

    if compression_level not in ['low', 'medium', 'high']:
        return "Invalid compression level", 400

    filename = secure_filename(pdf_file.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    pdf_file.save(filepath)

    try:
        original_size = os.path.getsize(filepath)
        logger.info(f"Original file size: {original_size / 1024:.2f} KB")

        output_filename = f"compressed_{uuid.uuid4().hex}.pdf"
        output_path = os.path.join(STATIC_FOLDER, output_filename)

        if is_ghostscript_installed():
            logger.info("Ghostscript detected, using it for compression")
            gs_quality = {'low': 'printer', 'medium': 'ebook', 'high': 'screen'}
            gs_setting = gs_quality[compression_level]
            
            gs_command = [
                'gs', '-sDEVICE=pdfwrite', '-dCompatibilityLevel=1.4',
                f'-dPDFSETTINGS=/{gs_setting}', '-dNOPAUSE', '-dQUIET', '-dBATCH',
                f'-sOutputFile={output_path}', filepath
            ]
            result = subprocess.run(gs_command, capture_output=True, text=True)
            if result.returncode != 0:
                logger.error(f"Ghostscript failed: {result.stderr}")
                raise Exception(f"Ghostscript compression failed: {result.stderr}")

        else:
            logger.warning("Ghostscript not found, falling back to image compression")
            reader = PdfReader(filepath)
            writer = PdfWriter()
            quality_map = {'low': 75, 'medium': 40, 'high': 10}
            jpeg_quality = quality_map[compression_level]
            image_count = 0

            for page in reader.pages:
                new_page = writer.add_page(page)
                if '/XObject' in page.get('/Resources', {}):
                    xobjects = page['/Resources']['/XObject'].get_object()
                    for obj in xobjects:
                        img_obj = xobjects[obj].get_object()
                        if isinstance(img_obj, dict) and img_obj.get('/Subtype') == '/Image' and img_obj.get('/Filter') in ['/DCTDecode', '/FlateDecode']:
                            try:
                                img_data = img_obj.get_data()
                                pil_img = Image.open(io.BytesIO(img_data))
                                if pil_img.mode != 'RGB':
                                    pil_img = pil_img.convert('RGB')

                                max_size = 1000 if compression_level == 'high' else 1500 if compression_level == 'medium' else 2000
                                if max(pil_img.size) > max_size:
                                    scale = max_size / max(pil_img.size)
                                    new_size = (int(pil_img.size[0] * scale), int(pil_img.size[1] * scale))
                                    pil_img = pil_img.resize(new_size, Image.Resampling.LANCZOS)

                                output = io.BytesIO()
                                pil_img.save(output, format='JPEG', quality=jpeg_quality, optimize=True)
                                compressed_data = output.getvalue()
                                
                                if len(compressed_data) < len(img_data):
                                    img_obj._data = compressed_data
                                    img_obj['/Length'] = len(img_obj._data)
                                    img_obj['/Filter'] = '/DCTDecode'
                                    img_obj['/ColorSpace'] = '/DeviceRGB'
                                    image_count += 1
                                else:
                                    logger.info(f"Skipping image {obj}: compressed size ({len(compressed_data)} bytes) not smaller than original ({len(img_data)} bytes)")
                            except Exception as e:
                                logger.warning(f"Failed to process image {obj}: {str(e)}")
                                continue

            writer.add_metadata(reader.metadata or {})
            logger.info(f"Processed {image_count} images for compression")
            with open(output_path, 'wb') as f:
                writer.write(f)

        compressed_size = os.path.getsize(output_path)
        logger.info(f"Compressed file size: {compressed_size / 1024:.2f} KB")
        reduction = (original_size - compressed_size) / original_size * 100 if original_size > 0 else 0
        logger.info(f"Size reduction: {reduction:.2f}%")

        log_conversion('compress-pdf', filename, output_filename, output_path)

        os.remove(filepath)
        return send_file(output_path, as_attachment=True, download_name='compressed.pdf')

    except Exception as e:
        os.remove(filepath)
        logger.error(f"Error during compression: {str(e)}")
        return f"Error during compression: {str(e)}", 500

@app.route('/split', methods=['POST'])
@login_required
def split_pdf():
    pdf_file = request.files.get('pdf')
    split_index = request.form.get('split_index')

    if not pdf_file or not split_index:
        return "Missing file or split index", 400

    try:
        split_index = int(split_index)
    except ValueError:
        return "Invalid split index", 400

    filename = secure_filename(pdf_file.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    pdf_file.save(filepath)

    reader = PdfReader(filepath)
    if split_index <= 0 or split_index >= len(reader.pages):
        os.remove(filepath)
        return "Split index out of range", 400

    writer1 = PdfWriter()
    writer2 = PdfWriter()

    for i, page in enumerate(reader.pages):
        if i < split_index:
            writer1.add_page(page)
        else:
            writer2.add_page(page)

    output_filename1 = f"split_part1_{uuid.uuid4().hex}.pdf"
    output_filename2 = f"split_part2_{uuid.uuid4().hex}.pdf"
    path1 = os.path.join(STATIC_FOLDER, output_filename1)
    path2 = os.path.join(STATIC_FOLDER, output_filename2)

    with open(path1, 'wb') as f:
        writer1.write(f)
    with open(path2, 'wb') as f:
        writer2.write(f)

    log_conversion('split-pdf', filename, f"{output_filename1}, {output_filename2}", f"{path1}, {path2}")

    os.remove(filepath)
    return jsonify({'part1': f'/static/{output_filename1}', 'part2': f'/static/{output_filename2}'})

# Document Screener Routes
@app.route('/analyze_document', methods=['POST'])
@login_required
def analyze_document():
    global current_document_text
    try:
        if 'docFile' not in request.files:
            return jsonify({'error': 'No document provided'}), 400
        
        doc_file = request.files['docFile']
        format_type = request.form.get('format', 'paragraph')
        
        temp_file_path = None
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{doc_file.filename.split('.')[-1]}") as temp_file:
            doc_file.save(temp_file)
            temp_file_path = temp_file.name
        
        file_extension = doc_file.filename.lower().split('.')[-1]
        if file_extension == 'pdf':
            text = extract_text_from_pdf(temp_file_path)
        elif file_extension == 'docx':
            text = extract_text_from_docx(temp_file_path)
        elif file_extension == 'txt':
            text = extract_text_from_txt(temp_file_path)
        else:
            os.remove(temp_file_path)
            return jsonify({'error': 'Unsupported file type. Use PDF, DOCX, or TXT.'}), 400
        
        os.remove(temp_file_path)
        
        if text.startswith('Error'):
            return jsonify({'error': text}), 500
        
        current_document_text = text
        
        analysis = analyze_text_with_openrouter(text, format_type)
        if analysis.startswith('Error'):
            return jsonify({'error': analysis}), 500
        
        analysis_filename = f"analysis_{uuid.uuid4().hex}.txt"
        analysis_path = os.path.join(OUTPUT_FOLDER, analysis_filename)
        with open(analysis_path, 'w', encoding='utf-8') as f:
            f.write(f"File: {doc_file.filename}\nFormat: {format_type}\nAnalysis:\n{analysis}\n\n")
        
        log_conversion('document-screener', doc_file.filename, analysis_filename, analysis_path)
        return jsonify({'analysis': analysis})
    except Exception as e:
        return jsonify({'error': f"Document analysis failed: {str(e)}"}), 500

@app.route('/chat', methods=['POST'])
@login_required
def chat():
    try:
        data = request.get_json()
        message = data.get('message', '')
        if not message:
            return jsonify({'error': 'No message provided'}), 400
        
        if not current_document_text:
            return jsonify({'error': 'No document uploaded. Please upload a document first.'}), 400
        
        response = chat_with_openrouter(message)
        if response.startswith('Error'):
            return jsonify({'error': response}), 500
        
        return jsonify({'response': response})
    except Exception as e:
        return jsonify({'error': f"Chat failed: {str(e)}"}), 500

# Plagiarism Scanner Route
@app.route('/check_plagiarism', methods=['POST'])
@login_required
def check_plagiarism():
    try:
        input_text = request.form.get('text', '').strip()
        if not input_text:
            return jsonify({'error': 'No text provided'}), 400

        if len(input_text) < 20:
            results_filename = f"plagiarism_results_{uuid.uuid4().hex}.json"
            results_path = os.path.join(OUTPUT_FOLDER, results_filename)
            results = [{"snippet": "", "similarity": "Input text too short to check plagiarism."}]
            with open(results_path, 'w', encoding='utf-8') as f:
                json.dump(results, f, indent=2)
            log_conversion('plagiarism-scanner', 'user_input.txt', results_filename, results_path)
            return jsonify({'results': results})

        query = input_text[:100]
        snippets = fetch_web_snippets(query)
        results = []
        for snippet in snippets:
            sim = call_openrouter_similarity(input_text, snippet)
            results.append({"snippet": snippet, "similarity": sim})

        results_filename = f"plagiarism_results_{uuid.uuid4().hex}.json"
        results_path = os.path.join(OUTPUT_FOLDER, results_filename)
        with open(results_path, 'w', encoding='utf-8') as f:
            json.dump(results, f, indent=2)

        log_conversion('plagiarism-scanner', 'user_input.txt', results_filename, results_path)
        return jsonify({'results': results})
    except Exception as e:
        logger.error(f"Plagiarism check failed: {str(e)}")
        return jsonify({'error': f"Plagiarism check failed: {str(e)}"}), 500

# Text to Speech Routes
@app.route('/generate_tts', methods=['POST'])
@login_required
def generate_tts():
    data = request.get_json()
    text = data.get('text', '').strip()
    
    if not text:
        return jsonify({'error': 'No text provided'}), 400

    try:
        filename = f"{uuid.uuid4()}.mp3"
        filepath = os.path.join(AUDIO_FOLDER, filename)
        tts = gTTS(text=text, lang='en')
        tts.save(filepath)
        log_conversion('text-to-speech', 'user_input.txt', filename, filepath)
        return jsonify({'audio_url': f'/download_audio/{filename}'})
    except Exception as e:
        return jsonify({'error': f"Text to speech conversion failed: {str(e)}"}), 500

@app.route('/download_audio/<filename>')
@login_required
def download_audio(filename):
    filepath = os.path.join(AUDIO_FOLDER, filename)
    if os.path.exists(filepath):
        return send_file(filepath, as_attachment=True)
    return jsonify({'error': 'File not found'}), 404

# Speech to Text Routes
@app.route('/save_transcript', methods=['POST'])
@login_required
def save_transcript():
    try:
        data = request.get_json()
        transcript = data.get('transcript', '')
        if not transcript:
            return jsonify({'error': 'No transcript provided'}), 400
        
        transcript_filename = f"transcript_{uuid.uuid4().hex}.txt"
        transcript_path = os.path.join(OUTPUT_FOLDER, transcript_filename)
        with open(transcript_path, 'w', encoding='utf-8') as f:
            f.write(transcript + '\n')
        log_conversion('speech-to-text', 'transcript.txt', transcript_filename, transcript_path)
        return jsonify({'message': 'Transcript saved successfully'})
    except Exception as e:
        return jsonify({'error': f"Transcript save failed: {str(e)}"}), 500

@app.route('/upload_audio', methods=['POST'])
@login_required
def upload_audio():
    try:
        if 'audioFile' not in request.files:
            return jsonify({'error': 'No audio file provided'}), 400
        
        audio_file = request.files['audioFile']
        temp_file_path = None
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as temp_file:
            audio_file.save(temp_file)
            temp_file_path = temp_file.name
        
        audio = AudioSegment.from_mp3(temp_file_path)
        wav_path = temp_file_path.replace('.mp3', '.wav')
        audio.export(wav_path, format='wav')
        
        with sr.AudioFile(wav_path) as source:
            audio_data = recognizer.record(source)
            try:
                transcript = recognizer.recognize_google(audio_data)
            except sr.UnknownValueError:
                transcript = "Could not understand audio"
            except sr.RequestError as e:
                transcript = f"Recognition error: {str(e)}"
        
        os.remove(temp_file_path)
        os.remove(wav_path)
        
        log_conversion('speech-to-text', audio_file.filename, 'transcript.txt', None)
        return jsonify({'transcript': transcript})
    except Exception as e:
        return jsonify({'error': f"Speech to text conversion failed: {str(e)}"}), 500

# AI PDF Editor Routes
@app.route('/analyze', methods=['POST'])
@login_required
def analyze():
    global latest_text
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    uploaded_file = request.files['file']
    if not uploaded_file.filename.lower().endswith('.pdf'):
        return jsonify({'error': 'Only PDF files are supported'}), 400

    filename = secure_filename(uploaded_file.filename)
    file_path = os.path.join(UPLOAD_FOLDER, filename)
    uploaded_file.save(file_path)

    try:
        extracted_text = extract_structured_text(file_path)
        if extracted_text.startswith('Error'):
            os.remove(file_path)
            return jsonify({'error': extracted_text}), 500
        if not extracted_text.strip():
            os.remove(file_path)
            return jsonify({'error': 'No text extracted from PDF. Ensure the PDF contains selectable text, not images.'}), 400
        latest_text = extracted_text

        prompt = (
            "You are an intelligent assistant analyzing a PDF document. Your task is to identify blank or unfilled fields such as 'Date: ____', 'Name: ________', 'Signature: [____]', or other placeholders (e.g., '________', '[____]', empty lines after labels). For each identified field, suggest a reasonable completion based on context (e.g., use today's date 'June 25, 2025' for date fields, 'John Doe' for name fields, 'Signature' for signature fields). Return the results in the following format:\n"
            "Identified Fields:\n"
            "- Field: [Description], Suggestion: [Suggested Value]\n"
            "If no blank fields are found, state: 'No blank or unfilled fields detected.'\n\n"
            f"Document Text:\n{extracted_text[:2000]}"  # Limit to avoid token overflow
        )

        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type": "application/json"
        }
        data = {
            "model": "mistralai/mixtral-8x7b-instruct",
            "messages": [{"role": "user", "content": prompt}],
            "max_tokens": 1000
        }

        try:
            response = retry_api_request(OPENROUTER_API_URL, headers, data)
            result = response.json()
            logger.debug(f"API response: {json.dumps(result, indent=2)}")
            suggestions = result.get('choices', [{}])[0].get('message', {}).get('content', '')
            if not suggestions:
                raise ValueError("Empty response content from API")
            log_conversion('ai-pdf-editor', filename, 'analysis.json', None)
            return jsonify({"text": extracted_text, "suggestions": suggestions})
        except requests.exceptions.RequestException as e:
            logger.error(f"OpenRouter API request failed: {str(e)}")
            return jsonify({'error': f"Failed to analyze document: API request error - {str(e)}"}), 500
        except ValueError as e:
            logger.error(f"OpenRouter API response error: {str(e)}")
            return jsonify({'error': f"Failed to analyze document: Invalid API response - {str(e)}"}), 500
    except Exception as e:
        logger.error(f"PDF analysis failed: {str(e)}")
        return jsonify({'error': f"Failed to analyze document: {str(e)}"}), 500
    finally:
        if os.path.exists(file_path):
            os.remove(file_path)

@app.route('/edit', methods=['POST'])
@login_required
def edit():
    try:
        updated_text = request.json.get('updated_text')
        if not updated_text:
            return jsonify({'error': 'No updated text provided'}), 400

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Arial", size=12)

        for line in updated_text.split('\n'):
            pdf.multi_cell(0, 10, line)

        output_filename = f"edited_{uuid.uuid4().hex}.pdf"
        output_path = os.path.join(UPLOAD_FOLDER, output_filename)
        pdf.output(output_path)

        log_conversion('ai-pdf-editor', 'user_input.txt', output_filename, output_path)
        return send_file(
            output_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name="edited_document.pdf"
        )
    except Exception as e:
        logger.error(f"PDF edit failed: {str(e)}")
        return jsonify({'error': f"Failed to edit document: {str(e)}"}), 500

@app.route('/fill_from_prompt', methods=['POST'])
@login_required
def fill_from_prompt():
    global latest_text
    try:
        user_prompt = request.json.get('user_prompt')
        if not user_prompt:
            return jsonify({'error': 'No user prompt provided'}), 400

        ai_instruction = (
            "Based on the following document text, a user wants to add or update content as follows:\n"
            f"Instruction: {user_prompt}\n\n"
            f"Document:\n{latest_text[:2000]}\n\n"
            "Provide the updated version of the document with the user request applied appropriately. "
            "Return only the revised document content."
        )

        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type": "application/json"
        }
        data = {
            "model": "mistralai/mixtral-8x7b-instruct",
            "messages": [{"role": "user", "content": ai_instruction}],
            "max_tokens": 1000
        }

        response = retry_api_request(OPENROUTER_API_URL, headers, data)
        result = response.json()
        updated_text = result.get('choices', [{}])[0].get('message', {}).get('content', '')
        if not updated_text:
            return jsonify({'error': 'Empty response from API'}), 500
        latest_text = updated_text
        log_conversion('ai-pdf-editor', 'user_prompt.txt', 'updated_text.txt', None)
        return jsonify({"updated_text": updated_text})
    except Exception as e:
        logger.error(f"Prompt-based edit failed: {str(e)}")
        return jsonify({'error': f"Prompt-based edit failed: {str(e)}"}), 500

# Text Summarizer Route
@app.route('/summarize', methods=['POST'])
@login_required
def summarize():
    try:
        data = request.get_json()
        text = data.get('text', '').strip()
        if not text:
            return jsonify({'error': 'No text provided'}), 400

        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type": "application/json"
        }

        payload = {
            "model": "mistralai/mixtral-8x7b-instruct",
            "messages": [
                {
                    "role": "user",
                    "content": (
                        "Summarize the following text in 3 sentences. "
                        "Only return the summary content. Do not add any introduction, title, or prefix:\n\n"
                        f"{text[:2000]}"
                    )
                }
            ],
            "max_tokens": 500
        }

        response = retry_api_request(OPENROUTER_API_URL, headers, payload)
        result = response.json()
        summary = result.get('choices', [{}])[0].get('message', {}).get('content', '').strip()
        if not summary:
            return jsonify({'error': 'Empty summary from API'}), 500

        summary_filename = f"summary_{uuid.uuid4().hex}.txt"
        summary_path = os.path.join(OUTPUT_FOLDER, summary_filename)
        with open(summary_path, 'w', encoding='utf-8') as f:
            f.write(summary)

        log_conversion('text-summarizer', 'user_input.txt', summary_filename, summary_path)
        return jsonify({'summary': summary})
    except Exception as e:
        logger.error(f"Text summarization failed: {str(e)}")
        return jsonify({'error': f"Text summarization failed: {str(e)}"}), 500

if __name__ == '__main__':
    app.run(debug=True)

# MAIN APP FOR TARS DOCSHIFT
# ALL MODULES WORKING