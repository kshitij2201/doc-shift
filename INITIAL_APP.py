from flask import Flask, request, send_file, jsonify, render_template, session, redirect, url_for
from functools import wraps
from PIL import Image
import os
import uuid
import sqlite3
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from datetime import datetime
from io import BytesIO
from pdf2image import convert_from_bytes
from docx import Document
from openpyxl import load_workbook
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from rembg import remove
import logging
import subprocess  # Added for Ghostscript in compress_pdf
import io

app = Flask(__name__)
app.secret_key = os.urandom(24)

# Configure logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

UPLOAD_FOLDER = 'Uploads'
CONVERTED_FOLDER = 'converted'
OUTPUT_FOLDER = 'converted'
STATIC_FOLDER = 'static'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(CONVERTED_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
os.makedirs(STATIC_FOLDER, exist_ok=True)

def init_db():
    conn = sqlite3.connect('file_conversion.db')
    c = conn.cursor()
    # Create the conversions table if it doesn't exist
    c.execute('''
        CREATE TABLE IF NOT EXISTS conversions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            conversion_type TEXT NOT NULL,
            original_filename TEXT NOT NULL,
            converted_filename TEXT NOT NULL,
            timestamp DATETIME DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    # Add file_path column if it doesn't exist
    try:
        c.execute('ALTER TABLE conversions ADD COLUMN file_path TEXT')
    except sqlite3.OperationalError:
        # Column already exists
        pass
    c.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT NOT NULL UNIQUE,
            password_hash TEXT NOT NULL
        )
    ''')
    try:
        default_username = 'admin'
        default_password = 'password'
        password_hash = generate_password_hash(default_password)
        c.execute('INSERT INTO users (username, password_hash) VALUES (?, ?)', (default_username, password_hash))
        conn.commit()
    except sqlite3.IntegrityError:
        pass
    conn.close()

init_db()

def log_conversion(conversion_type, original_filename, converted_filename, file_path=None):
    conn = sqlite3.connect('file_conversion.db')
    c = conn.cursor()
    c.execute('''
        INSERT INTO conversions (conversion_type, original_filename, converted_filename, file_path)
        VALUES (?, ?, ?, ?)
    ''', (conversion_type, original_filename, converted_filename, file_path))
    conn.commit()
    conn.close()

def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'logged_in' not in session:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

# Helper function for Ghostscript check (used in compress_pdf)
def is_ghostscript_installed():
    try:
        subprocess.run(['gs', '--version'], stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
        return True
    except (subprocess.CalledProcessError, FileNotFoundError):
        return False

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        conn = sqlite3.connect('file_conversion.db')
        c = conn.cursor()
        c.execute('SELECT password_hash FROM users WHERE username = ?', (username,))
        user = c.fetchone()
        conn.close()
        if user and check_password_hash(user[0], password):
            session['logged_in'] = True
            return redirect(url_for('index'))
        else:
            return render_template('login.html', error='Invalid username or password')
    return render_template('login.html', error=None)

@app.route('/logout')
@login_required
def logout():
    session.pop('logged_in', None)
    return redirect(url_for('login'))

@app.route('/')
@login_required
def index():
    return render_template('index.html')

@app.route('/image-to-pdf')
@login_required
def image_to_pdf_page():
    return render_template('image_to_pdf.html')

@app.route('/pdf-to-image')
@login_required
def pdf_to_image_page():
    return render_template('pdf_to_image.html')

@app.route('/merge-pdfs')
@login_required
def merge_pdfs_page():
    return render_template('merge_pdfs.html')

@app.route('/word-to-pdf')
@login_required
def word_to_pdf_page():
    return render_template('word_to_pdf.html')

@app.route('/excel-to-pdf')
@login_required
def excel_to_pdf_page():
    return render_template('excel_to_pdf.html')

@app.route('/pdf-to-ppt')
@login_required
def pdf_to_ppt_page():
    return render_template('pdf_to_ppt.html')

@app.route('/bg-remover')
@login_required
def bg_remover_page():
    return render_template('bg_remover.html')

@app.route('/admin-logs')
@login_required
def logs_page():
    return render_template('logs.html')

@app.route('/compress-pdf')
@login_required
def compress_pdf_page():
    return render_template('compress_pdf.html')

@app.route('/convert/image-to-pdf', methods=['POST'])
@login_required
def convert_image_to_pdf():
    if 'images' not in request.files:
        return jsonify({'error': 'No images provided'}), 400

    files = request.files.getlist('images')
    image_list = []

    for file in files:
        try:
            image = Image.open(file.stream)
            if image.mode != 'RGB':
                image = image.convert('RGB')
            image_list.append(image)
        except Exception as e:
            return jsonify({'error': f'Failed to read image: {str(e)}'}), 500

    if not image_list:
        return jsonify({'error': 'No valid images found'}), 400

    output_filename = f"{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(CONVERTED_FOLDER, output_filename)

    try:
        image_list[0].save(output_path, save_all=True, append_images=image_list[1:], format='PDF')
        log_conversion('image-to-pdf', files[0].filename, output_filename, output_path)
    except Exception as e:
        return jsonify({'error': f'PDF conversion failed: {str(e)}'}), 500

    return send_file(output_path, as_attachment=True, download_name="converted.pdf", mimetype='application/pdf')

@app.route('/convert/pdf-to-image', methods=['POST'])
@login_required
def convert_pdf_to_image():
    pdf_file = request.files.get('pdf')
    if not pdf_file:
        return "No PDF uploaded", 400

    try:
        images = convert_from_bytes(pdf_file.read(), fmt='png', single_file=True)
        output_filename = f"{uuid.uuid4().hex}.png"
        output_path = os.path.join(CONVERTED_FOLDER, output_filename)
        images[0].save(output_path, format='PNG')

        log_conversion('pdf-to-image', pdf_file.filename, output_filename, output_path)

        return send_file(output_path, as_attachment=True, download_name="converted.png", mimetype='image/png')
    except Exception as e:
        return f"Error during conversion: {str(e)}", 500

@app.route('/merge/pdfs', methods=['POST'])
@login_required
def merge_pdfs():
    if 'pdfs' not in request.files:
        return jsonify({'error': 'No PDF files provided'}), 400

    files = request.files.getlist('pdfs')
    merger = PdfMerger()

    try:
        for file in files:
            merger.append(file)

        output_filename = f"merged_{uuid.uuid4().hex}.pdf"
        output_path = os.path.join(CONVERTED_FOLDER, output_filename)
        merger.write(output_path)
        merger.close()

        log_conversion('merge-pdfs', files[0].filename, output_filename, output_path)
        return send_file(output_path, as_attachment=True, download_name='merged.pdf', mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Merge failed: {str(e)}'}), 500

@app.route('/admin/logs')
@login_required
def view_logs():
    conn = sqlite3.connect('file_conversion.db')
    c = conn.cursor()
    c.execute('SELECT * FROM conversions ORDER BY timestamp DESC')
    rows = c.fetchall()
    conn.close()
    return jsonify(rows)

@app.route('/download')
@login_required
def download():
    file_path = request.args.get('file_path')
    file_name = request.args.get('file_name')
    mime_type = request.args.get('mime_type', 'application/octet-stream')
    if not os.path.exists(file_path):
        return jsonify({'error': 'File not found on server'}), 404
    return send_file(file_path, as_attachment=True, download_name=file_name, mimetype=mime_type)

@app.route('/convert_word_to_pdf', methods=['POST'])
@login_required
def convert_word_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded!'}), 400

    file = request.files['file']
    if not file.filename.endswith('.docx'):
        return jsonify({'error': 'Only .docx files are supported'}), 400

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_filename = f"{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(OUTPUT_FOLDER, output_filename)

    try:
        file.save(input_path)
        doc = Document(input_path)
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        pdf_io = BytesIO()
        c = canvas.Canvas(pdf_io, pagesize=letter)
        y = 750
        for line in text:
            if y < 50:
                c.showPage()
                y = 750
            c.drawString(50, y, line[:100])
            y -= 15
        c.save()
        pdf_io.seek(0)
        log_conversion('word-to-pdf', file.filename, output_filename, output_path)
        with open(output_path, 'wb') as f:
            f.write(pdf_io.getvalue())
        return send_file(output_path, as_attachment=True, download_name=file.filename.replace('.docx', '.pdf'), mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Word to PDF conversion failed: {str(e)}'}), 500
    finally:
        if os.path.exists(input_path):
            os.remove(input_path)

@app.route('/convert_excel_to_pdf', methods=['POST'])
@login_required
def convert_excel_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded!'}), 400

    file = request.files['file']
    if not file.filename.endswith('.xlsx'):
        return jsonify({'error': 'Only .xlsx files are supported'}), 400

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_filename = f"{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(OUTPUT_FOLDER, output_filename)

    try:
        file.save(input_path)
        wb = load_workbook(input_path)
        sheet = wb.active
        text = []
        for row in sheet.rows:
            row_text = [str(cell.value) if cell.value is not None else '' for cell in row]
            text.append(' | '.join(row_text))
        pdf_io = BytesIO()
        c = canvas.Canvas(pdf_io, pagesize=letter)
        y = 750
        for line in text:
            if y < 50:
                c.showPage()
                y = 750
            c.drawString(50, y, line[:100])
            y -= 15
        c.save()
        pdf_io.seek(0)
        log_conversion('excel-to-pdf', file.filename, output_filename, output_path)
        with open(output_path, 'wb') as f:
            f.write(pdf_io.getvalue())
        return send_file(output_path, as_attachment=True, download_name=file.filename.replace('.xlsx', '.pdf'), mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Excel to PDF conversion failed: {str(e)}'}), 500
    finally:
        if os.path.exists(input_path):
            os.remove(input_path)

@app.route('/convert_pdf_to_ppt', methods=['POST'])
@login_required
def convert_pdf_to_ppt():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded!'}), 400

    file = request.files['file']
    if not file.filename.endswith('.pdf'):
        return jsonify({'error': 'Only .pdf files are supported'}), 400

    output_filename = f"{uuid.uuid4().hex}.pptx"
    output_path = os.path.join(OUTPUT_FOLDER, output_filename)

    try:
        pdf_reader = PdfReader(file.stream)
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content layout

        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text = page.extract_text() or f"Page {page_num + 1}"
            lines = text.split('\n')
            filtered_lines = [line.strip() for line in lines if line.strip()]
            slide_text = '\n'.join(filtered_lines[:10])  # Limit to 10 lines to avoid overflow

            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = f"Slide {page_num + 1}"
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.text = slide_text
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        ppt_io = BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)

        log_conversion('pdf-to-ppt', file.filename, output_filename, output_path)
        with open(output_path, 'wb') as f:
            f.write(ppt_io.getvalue())
        return send_file(output_path, as_attachment=True, download_name=file.filename.replace('.pdf', '.pptx'), mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    except Exception as e:
        return jsonify({'error': f'PDF to PPT conversion failed: {str(e)}'}), 500

@app.route('/remove_background', methods=['POST'])
@login_required
def remove_background():
    logger.debug("Received request at /remove_background")
    
    if 'image' not in request.files:
        logger.error("No image uploaded")
        return jsonify({'error': 'No image uploaded!'}), 400

    file = request.files['image']
    valid_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp')
    if not file.filename.lower().endswith(valid_extensions):
        logger.error(f"Unsupported file extension: {file.filename}")
        return jsonify({'error': 'Only JPG, PNG, GIF, BMP, TIFF, and WEBP files are supported'}), 400

    if file.content_length > 10 * 1024 * 1024:
        logger.error("File size exceeds 10MB limit")
        return jsonify({'error': 'File size exceeds 10MB limit'}), 400

    try:
        logger.debug("Opening image with PIL")
        image = Image.open(file.stream)
        logger.debug(f"Image size: {image.size}")
        
        if image.size[0] > 4096 or image.size[1] > 4096:
            logger.error("Image resolution exceeds 4096x4096 limit")
            return jsonify({'error': 'Image resolution exceeds 4096x4096 limit'}), 400

        if file.filename.lower().endswith('.gif'):
            logger.debug("Processing GIF, using first frame")
            image.seek(0)

        if image.mode != 'RGBA':
            logger.debug("Converting image to RGBA mode")
            image = image.convert('RGBA')

        logger.debug("Starting background removal with rembg")
        output_image = remove(image)
        logger.debug("Background removal completed")

        output_filename = f"bg_removed_{uuid.uuid4().hex}.png"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        logger.debug(f"Saving output to disk: {output_path}")
        output_image.save(output_path, format='PNG')

        log_conversion('background-remover', file.filename, output_filename, output_path)

        logger.debug("Sending processed image back to client")
        return send_file(output_path, as_attachment=True, download_name='background_removed.png', mimetype='image/png')
    except Exception as e:
        logger.error(f"Background removal failed: {str(e)}", exc_info=True)
        return jsonify({'error': f'Background removal failed: {str(e)}'}), 500

@app.route('/remove-pages-ui')
@login_required
def remove_pages_ui():
    return render_template('remove_page.html')

@app.route('/get-page-count', methods=['POST'])
@login_required
def get_page_count():
    pdf = request.files.get('pdf')
    if not pdf:
        return jsonify({'error': 'No PDF uploaded'}), 400

    filename = secure_filename(pdf.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    pdf.save(filepath)

    try:
        reader = PdfReader(filepath)
        count = len(reader.pages)
        os.remove(filepath)
        return jsonify({'page_count': count})
    except Exception as e:
        os.remove(filepath)
        return jsonify({'error': str(e)}), 500

@app.route('/remove-pages', methods=['POST'])
@login_required
def remove_pages():
    pdf = request.files.get('pdf')
    removed_pages = request.form.get('removed_pages', '')
    total_pages = int(request.form.get('page_count', 0))

    if not pdf:
        return "No PDF uploaded", 400

    try:
        remove_indices = [int(i) for i in removed_pages.split(',') if i.strip().isdigit()]
    except ValueError:
        return "Invalid page indices", 400

    filename = secure_filename(pdf.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    pdf.save(filepath)

    reader = PdfReader(filepath)
    writer = PdfWriter()

    for i in range(len(reader.pages)):
        if i not in remove_indices:
            writer.add_page(reader.pages[i])

    output_filename = f"removed_{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(STATIC_FOLDER, output_filename)
    with open(output_path, 'wb') as f:
        writer.write(f)

    os.remove(filepath)
    log_conversion('remove-pages', filename, output_filename, output_path)
    return send_file(output_path, as_attachment=True, download_name='removed_pages.pdf')

@app.route('/compress', methods=['POST'])
@login_required
def compress_pdf():
    pdf_file = request.files.get('pdf')
    compression_level = request.form.get('compression_level')

    if not pdf_file:
        return "No PDF uploaded", 400

    if compression_level not in ['low', 'medium', 'high']:
        return "Invalid compression level", 400

    filename = secure_filename(pdf_file.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    pdf_file.save(filepath)

    try:
        # Get original file size for logging
        original_size = os.path.getsize(filepath)
        logger.info(f"Original file size: {original_size / 1024:.2f} KB")

        output_filename = f"compressed_{uuid.uuid4().hex}.pdf"  # Added unique filename
        output_path = os.path.join(STATIC_FOLDER, output_filename)

        # Try Ghostscript compression if available
        if is_ghostscript_installed():
            logger.info("Ghostscript detected, using it for compression")
            # Map compression levels to Ghostscript settings
            gs_quality = {'low': 'printer', 'medium': 'ebook', 'high': 'screen'}
            gs_setting = gs_quality[compression_level]
            
            # Run Ghostscript command
            gs_command = [
                'gs', '-sDEVICE=pdfwrite', '-dCompatibilityLevel=1.4',
                f'-dPDFSETTINGS=/{gs_setting}', '-dNOPAUSE', '-dQUIET', '-dBATCH',
                f'-sOutputFile={output_path}', filepath
            ]
            result = subprocess.run(gs_command, capture_output=True, text=True)
            if result.returncode != 0:
                logger.error(f"Ghostscript failed: {result.stderr}")
                raise Exception(f"Ghostscript compression failed: {result.stderr}")

        else:
            logger.warning("Ghostscript not found, falling back to image compression")
            # Fallback to image compression with pypdf
            reader = PdfReader(filepath)
            writer = PdfWriter()
            quality_map = {'low': 75, 'medium': 40, 'high': 10}
            jpeg_quality = quality_map[compression_level]
            image_count = 0

            for page in reader.pages:
                new_page = writer.add_page(page)
                if '/XObject' in page.get('/Resources', {}):
                    xobjects = page['/Resources']['/XObject'].get_object()
                    for obj in xobjects:
                        img_obj = xobjects[obj].get_object()
                        if isinstance(img_obj, dict) and img_obj.get('/Subtype') == '/Image' and img_obj.get('/Filter') in ['/DCTDecode', '/FlateDecode']:
                            try:
                                img_data = img_obj.get_data()
                                pil_img = Image.open(io.BytesIO(img_data))
                                if pil_img.mode != 'RGB':
                                    pil_img = pil_img.convert('RGB')

                                max_size = 1000 if compression_level == 'high' else 1500 if compression_level == 'medium' else 2000
                                if max(pil_img.size) > max_size:
                                    scale = max_size / max(pil_img.size)
                                    new_size = (int(pil_img.size[0] * scale), int(pil_img.size[1] * scale))
                                    pil_img = pil_img.resize(new_size, Image.Resampling.LANCZOS)

                                output = io.BytesIO()
                                pil_img.save(output, format='JPEG', quality=jpeg_quality, optimize=True)
                                compressed_data = output.getvalue()
                                
                                if len(compressed_data) < len(img_data):
                                    img_obj._data = compressed_data
                                    img_obj['/Length'] = len(img_obj._data)
                                    img_obj['/Filter'] = '/DCTDecode'
                                    img_obj['/ColorSpace'] = '/DeviceRGB'
                                    image_count += 1
                                else:
                                    logger.info(f"Skipping image {obj}: compressed size ({len(compressed_data)} bytes) not smaller than original ({len(img_data)} bytes)")
                            except Exception as e:
                                logger.warning(f"Failed to process image {obj}: {str(e)}")
                                continue

            writer.add_metadata(reader.metadata or {})
            logger.info(f"Processed {image_count} images for compression")
            with open(output_path, 'wb') as f:
                writer.write(f)

        # Log compressed file size
        compressed_size = os.path.getsize(output_path)
        logger.info(f"Compressed file size: {compressed_size / 1024:.2f} KB")
        reduction = (original_size - compressed_size) / original_size * 100 if original_size > 0 else 0
        logger.info(f"Size reduction: {reduction:.2f}%")

        # Log the conversion
        log_conversion('compress-pdf', filename, output_filename, output_path)

        os.remove(filepath)
        return send_file(output_path, as_attachment=True, download_name='compressed.pdf')

    except Exception as e:
        os.remove(filepath)
        logger.error(f"Error during compression: {str(e)}")
        return f"Error during compression: {str(e)}", 500

@app.route('/split-pdf')
@login_required
def split_pdf_page():
    return render_template('split_pdf.html')

@app.route('/split', methods=['POST'])
@login_required
def split_pdf():
    pdf_file = request.files.get('pdf')
    split_index = request.form.get('split_index')

    if not pdf_file or not split_index:
        return "Missing file or split index", 400

    try:
        split_index = int(split_index)
    except ValueError:
        return "Invalid split index", 400

    filename = secure_filename(pdf_file.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    pdf_file.save(filepath)

    reader = PdfReader(filepath)
    if split_index <= 0 or split_index >= len(reader.pages):
        os.remove(filepath)
        return "Split index out of range", 400

    writer1 = PdfWriter()
    writer2 = PdfWriter()

    for i, page in enumerate(reader.pages):
        if i < split_index:
            writer1.add_page(page)
        else:
            writer2.add_page(page)

    output_filename1 = f"split_part1_{uuid.uuid4().hex}.pdf"
    output_filename2 = f"split_part2_{uuid.uuid4().hex}.pdf"
    path1 = os.path.join(STATIC_FOLDER, output_filename1)
    path2 = os.path.join(STATIC_FOLDER, output_filename2)

    with open(path1, 'wb') as f:
        writer1.write(f)
    with open(path2, 'wb') as f:
        writer2.write(f)

    log_conversion('split-pdf', filename, f"{output_filename1}, {output_filename2}", f"{path1}, {path2}")

    os.remove(filepath)
    return jsonify({'part1': f'/static/{output_filename1}', 'part2': f'/static/{output_filename2}'})

if __name__ == '__main__':
    app.run(debug=True)



  # THIS IS AN INITIAL APP WITH LESS MODULES  